"""Document parsing service for Memora."""

import logging
import os
import tempfile
import time
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from uuid import uuid4

import aiofiles
from fastapi import UploadFile

from app.config import settings
from app.core.ingestion.chunker import TextChunker
from app.models.ingest import (
    ChunkingStrategy,
    DocumentChunk,
    DocumentType,
    ProcessedDocument,
)

logger = logging.getLogger(__name__)


class DocumentParser:
    """Service for parsing various document types."""

    # Supported file extensions
    SUPPORTED_EXTENSIONS = {
        ".pdf": DocumentType.PDF,
        ".docx": DocumentType.DOCX,
        ".doc": DocumentType.DOCX,
        ".pptx": DocumentType.PPTX,
        ".ppt": DocumentType.PPTX,
        ".xlsx": DocumentType.XLSX,
        ".xls": DocumentType.XLSX,
        ".txt": DocumentType.TXT,
        ".md": DocumentType.MD,
        ".html": DocumentType.HTML,
        ".htm": DocumentType.HTML,
        ".json": DocumentType.JSON,
        ".csv": DocumentType.CSV,
    }

    def __init__(self):
        """Initialize the document parser."""
        self.chunker = TextChunker()
        self._ensure_upload_dir()

    def _ensure_upload_dir(self) -> None:
        """Ensure upload directory exists."""
        os.makedirs(settings.upload_dir, exist_ok=True)

    def get_document_type(self, filename: str) -> Optional[DocumentType]:
        """Determine document type from filename."""
        ext = Path(filename).suffix.lower()
        return self.SUPPORTED_EXTENSIONS.get(ext)

    async def parse_file(
        self,
        file: UploadFile,
        chunking_strategy: ChunkingStrategy = ChunkingStrategy.SEMANTIC,
        chunk_size: int = 512,
        chunk_overlap: int = 50,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> ProcessedDocument:
        """
        Parse an uploaded file and extract chunks.
        
        Args:
            file: Uploaded file
            chunking_strategy: How to split the document
            chunk_size: Target chunk size
            chunk_overlap: Overlap between chunks
            metadata: Additional metadata
            
        Returns:
            ProcessedDocument with extracted chunks
        """
        start_time = time.time()
        
        # Validate file type
        doc_type = self.get_document_type(file.filename or "unknown.txt")
        if doc_type is None:
            raise ValueError(f"Unsupported file type: {file.filename}")

        # Save file temporarily
        temp_path = await self._save_temp_file(file)
        
        try:
            # Extract text based on document type
            text, extracted_metadata = await self._extract_content(temp_path, doc_type)
            
            # Configure chunker
            self.chunker = TextChunker(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                strategy=chunking_strategy,
            )
            
            # Merge metadata
            full_metadata = {**(metadata or {}), **extracted_metadata}
            
            # Chunk the text
            chunks = self.chunker.chunk_text(text, full_metadata)
            
            processing_time = (time.time() - start_time) * 1000
            
            return ProcessedDocument(
                document_id=uuid4(),
                filename=file.filename,
                document_type=doc_type,
                title=extracted_metadata.get("title", file.filename),
                chunks=chunks,
                total_chunks=len(chunks),
                total_pages=extracted_metadata.get("total_pages"),
                total_characters=len(text),
                extracted_metadata=extracted_metadata,
                processing_time_ms=processing_time,
            )
            
        finally:
            # Clean up temp file
            if os.path.exists(temp_path):
                os.remove(temp_path)

    async def parse_text(
        self,
        content: str,
        title: Optional[str] = None,
        chunking_strategy: ChunkingStrategy = ChunkingStrategy.SEMANTIC,
        chunk_size: int = 512,
        chunk_overlap: int = 50,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> ProcessedDocument:
        """
        Parse raw text content.
        
        Args:
            content: Raw text content
            title: Optional title
            chunking_strategy: How to split the content
            chunk_size: Target chunk size
            chunk_overlap: Overlap between chunks
            metadata: Additional metadata
            
        Returns:
            ProcessedDocument with extracted chunks
        """
        start_time = time.time()
        
        # Configure chunker
        self.chunker = TextChunker(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            strategy=chunking_strategy,
        )
        
        # Chunk the text
        chunks = self.chunker.chunk_text(content, metadata or {})
        
        processing_time = (time.time() - start_time) * 1000
        
        return ProcessedDocument(
            document_id=uuid4(),
            filename=None,
            document_type=DocumentType.TXT,
            title=title,
            chunks=chunks,
            total_chunks=len(chunks),
            total_pages=None,
            total_characters=len(content),
            extracted_metadata=metadata or {},
            processing_time_ms=processing_time,
        )

    async def _save_temp_file(self, file: UploadFile) -> str:
        """Save uploaded file to temporary location."""
        suffix = Path(file.filename or "file.tmp").suffix
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            content = await file.read()
            tmp.write(content)
            return tmp.name

    async def _extract_content(
        self, 
        file_path: str, 
        doc_type: DocumentType
    ) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text content from document.
        
        Returns:
            Tuple of (text_content, metadata)
        """
        extractors = {
            DocumentType.PDF: self._extract_pdf,
            DocumentType.DOCX: self._extract_docx,
            DocumentType.PPTX: self._extract_pptx,
            DocumentType.XLSX: self._extract_xlsx,
            DocumentType.TXT: self._extract_text,
            DocumentType.MD: self._extract_text,
            DocumentType.HTML: self._extract_html,
            DocumentType.JSON: self._extract_json,
            DocumentType.CSV: self._extract_csv,
        }
        
        extractor = extractors.get(doc_type, self._extract_text)
        return await extractor(file_path)

    async def _extract_pdf(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from PDF using pypdf."""
        try:
            from pypdf import PdfReader
            
            reader = PdfReader(file_path)
            
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text() or ""
                text_parts.append(page_text)
            
            text = "\n\n".join(text_parts)
            
            metadata = {
                "total_pages": len(reader.pages),
            }
            
            # Extract PDF metadata if available
            if reader.metadata:
                if reader.metadata.title:
                    metadata["title"] = reader.metadata.title
                if reader.metadata.author:
                    metadata["author"] = reader.metadata.author
                if reader.metadata.subject:
                    metadata["subject"] = reader.metadata.subject
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            raise ValueError(f"Failed to extract PDF content: {e}")

    async def _extract_docx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from DOCX."""
        try:
            from docx import Document
            
            doc = Document(file_path)
            
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                    if row_text:
                        text_parts.append(row_text)
            
            text = "\n\n".join(text_parts)
            
            metadata = {}
            if doc.core_properties.title:
                metadata["title"] = doc.core_properties.title
            if doc.core_properties.author:
                metadata["author"] = doc.core_properties.author
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise ValueError(f"Failed to extract DOCX content: {e}")

    async def _extract_pptx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from PowerPoint."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
            
            text = "\n\n".join(text_parts)
            
            metadata = {
                "total_pages": len(prs.slides),
            }
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise ValueError(f"Failed to extract PPTX content: {e}")

    async def _extract_xlsx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from Excel."""
        try:
            import pandas as pd
            
            excel_file = pd.ExcelFile(file_path)
            
            text_parts = []
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                sheet_text = f"[Sheet: {sheet_name}]\n{df.to_string(index=False)}"
                text_parts.append(sheet_text)
            
            text = "\n\n".join(text_parts)
            
            metadata = {
                "sheets": excel_file.sheet_names,
                "total_sheets": len(excel_file.sheet_names),
            }
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            raise ValueError(f"Failed to extract XLSX content: {e}")

    async def _extract_text(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from plain text file."""
        async with aiofiles.open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = await f.read()
        return text, {}

    async def _extract_html(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from HTML."""
        try:
            from bs4 import BeautifulSoup
            
            async with aiofiles.open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = await f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for element in soup(['script', 'style', 'meta', 'link']):
                element.decompose()
            
            text = soup.get_text(separator='\n', strip=True)
            
            metadata = {}
            title_tag = soup.find('title')
            if title_tag:
                metadata["title"] = title_tag.get_text()
            
            return text, metadata
            
        except ImportError:
            # Fallback if beautifulsoup not available
            return await self._extract_text(file_path)

    async def _extract_json(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from JSON."""
        import json
        
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            content = await f.read()
        
        data = json.loads(content)
        
        # Convert JSON to readable text
        text = json.dumps(data, indent=2)
        
        return text, {"type": "json"}

    async def _extract_csv(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract content from CSV."""
        try:
            import pandas as pd
            
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
            
            metadata = {
                "columns": list(df.columns),
                "rows": len(df),
            }
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"CSV extraction failed: {e}")
            # Fallback to raw text
            return await self._extract_text(file_path)


# Global parser instance
document_parser = DocumentParser()
